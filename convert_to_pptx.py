import json
import os
import re
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def load_songs():
    with open('data/songs.json', 'r', encoding='utf-8') as f:
        return json.load(f)

def create_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    
    # Calculate center position
    # Slide width is 10 inches, height is 5.625 inches
    width = Inches(8)  # text box width
    height = Inches(3)  # text box height
    left = (prs.slide_width - width) / 2
    top = (prs.slide_height - height) / 2
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Center align the text
    p = tf.add_paragraph()
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    
    return slide

def add_amin_to_slide(slide):
    # Add "Amin!" to bottom right
    left = Inches(8)
    top = Inches(4.5)  # Adjusted to be higher up from the bottom
    width = Inches(2)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Amin!"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(24)

def create_pptx_for_song(song):
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add slides for each verse/stanza
    slides = song['slides']
    for i, slide_text in enumerate(slides):
        slide = create_slide(prs, slide_text)
        # Add "Amin!" to the last slide with content
        if i == len(slides) - 1:
            add_amin_to_slide(slide)
    
    # Create directory structure
    output_dir = os.path.join('data', 'pptx', song['type'])
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    # Clean up title by removing leading and trailing non-letters except '?' and spaces, keeping diacritics
    clean_title = re.sub(r'^[^a-zA-ZăâîșțĂÂÎȘȚ]+', '', song['title']).strip()  # Remove leading non-letters
    clean_title = re.sub(r'[^a-zA-ZăâîșțĂÂÎȘȚ\?]+$', '', clean_title).strip()  # Remove trailing non-letters except '?'
    filename = f"{clean_title}.pptx"
    # Remove any remaining invalid characters from filename
    filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_', '.'))
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)

def generate_detailed_html(web_dir, songs):
    # Group songs by type
    songs_by_type = {}
    for song in songs:
        if song['type'] not in songs_by_type:
            songs_by_type[song['type']] = []
        songs_by_type[song['type']].append(song)
    
    html_content = f'''<!DOCTYPE html>
<html lang="ro">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cântări Creștine Videoproiector - Listă completă</title>
    <style>
        body {{
            font-family: system-ui, -apple-system, sans-serif;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
            line-height: 1.6;
        }}
        h1, h2 {{
            color: #2c3e50;
            text-align: center;
        }}
        h2 {{
            font-size: 1rem;
            margin-top: -12px;
            margin-bottom: 29px;
        }}
        .search-container {{
            margin: 20px auto;
            max-width: 600px;
            text-align: center;
        }}
        #searchInput {{
            width: 100%;
            padding: 10px;
            font-size: 16px;
            border: 2px solid #3498db;
            border-radius: 5px;
            margin-bottom: 20px;
        }}
        .category {{
            margin-bottom: 40px;
            padding: 20px;
            background: #f8f9fa;
            border-radius: 10px;
        }}
        .category h3 {{
            color: #2c3e50;
            margin-top: 0;
        }}
        .songs-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(250px, 1fr));
            gap: 10px;
        }}
        .song-link {{
            display: block;
            padding: 10px;
            background: #3498db;
            color: white;
            text-decoration: none;
            border-radius: 5px;
            transition: background 0.3s;
        }}
        .song-link:hover {{
            background: #2980b9;
        }}
        .back-link {{
            display: inline-block;
            margin-bottom: 20px;
            color: #3498db;
            text-decoration: none;
        }}
        .back-link:hover {{
            text-decoration: underline;
        }}
        .hidden {{
            display: none;
        }}
    </style>
</head>
<body>
    <a href="index.html" class="back-link">← Înapoi la pagina de descărcare</a>
    <h1>Cântări Creștine Videoproiector - Listă completă</h1>
    <h2>Cântările au fost preluate de pe <a href="https://cantari-crestine.com/" target="_blank">cantari-crestine.com</a></h2>
    
    <div class="search-container">
        <input type="text" id="searchInput" placeholder="Caută o cântare...">
    </div>

    <div id="songList">'''

    # Add each category and its songs
    for category, category_songs in songs_by_type.items():
        html_content += f'''
        <div class="category" data-category="{category}">
            <h3>{category}</h3>
            <div class="songs-grid">'''
        
        for song in sorted(category_songs, key=lambda x: x['title']):
            clean_title = re.sub(r'^[^a-zA-ZăâîșțĂÂÎȘȚ]+', '', song['title']).strip()
            clean_title = re.sub(r'[^a-zA-ZăâîșțĂÂÎȘȚ\?]+$', '', clean_title).strip()
            filename = f"{clean_title}.pptx"
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_', '.'))
            file_path = f"{song['type_abvr']}/{filename}"
            
            html_content += f'''
                <a href="{file_path}" class="song-link" download title="{song['title']}">{song['title']}</a>'''
        
        html_content += '''
            </div>
        </div>'''

    html_content += '''
    </div>

    <script>
        document.getElementById('searchInput').addEventListener('input', function(e) {
            const searchText = e.target.value.toLowerCase();
            const songLinks = document.getElementsByClassName('song-link');
            const categories = document.getElementsByClassName('category');
            
            for (let category of categories) {
                let hasVisibleSongs = false;
                const songs = category.getElementsByClassName('song-link');
                
                for (let song of songs) {
                    const title = song.textContent.toLowerCase();
                    if (title.includes(searchText)) {
                        song.style.display = 'block';
                        hasVisibleSongs = true;
                    } else {
                        song.style.display = 'none';
                    }
                }
                
                category.style.display = hasVisibleSongs ? 'block' : 'none';
            }
        });
    </script>
</body>
</html>'''
    
    with open(os.path.join(web_dir, 'cantari-crestine-videoproiector.html'), 'w', encoding='utf-8') as f:
        f.write(html_content)

def generate_index_html(web_dir, zip_files):
    links_html = []
    base_url = "https://raw.githubusercontent.com/radio-crestin/cantari-crestine-videoproiector/main/data/pptx"
    
    # Sort and process zip files, putting the main archive first
    main_zip = 'Toate cantarile crestine videoproiector.zip'
    sorted_files = sorted([f for f in zip_files if f != main_zip])
    if main_zip in zip_files:
        # Add main zip file with special styling
        url = f"{base_url}/{main_zip}"
        links_html.append(f'<a href="{url}" class="download-link main-download" target="_blank" download>{main_zip.replace(".zip", "")} (arhiva .zip)</a>')
    
    # Add other zip files
    for zip_file in sorted_files:
        name = zip_file.replace('.zip', '')
        url = f"{base_url}/{zip_file}"
        links_html.append(f'<a href="{url}" class="download-link" target="_blank" download>{name} (arhiva .zip)</a>')
    
    html_content = f'''<!DOCTYPE html>
<html lang="ro">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cântări Creștine Videoproiector Download</title>
    <style>
        body {{
            font-family: system-ui, -apple-system, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
            line-height: 1.6;
        }}
        h1 {{
            color: #2c3e50;
            text-align: center;
        }}
        h2 {{
            color: #2c3e50;
            text-align: center;
            font-size: 1rem;
            margin-top: -12px;
            margin-bottom: 29px;
        }}
        .downloads {{
            display: flex;
            flex-direction: column;
            gap: 10px;
        }}
        .download-link {{
            display: block;
            padding: 10px 15px;
            background: #3498db;
            color: white;
            text-decoration: none;
            border-radius: 5px;
            transition: background 0.3s;
            margin-bottom: 5px;
        }}
        .download-link:hover {{
            background: #2980b9;
        }}
        .main-download {{
            background: #e74c3c;
            font-size: 1.2em;
            margin-bottom: 20px;
        }}
        .main-download:hover {{
            background: #c0392b;
        }}
    </style>
</head>
<body>
    <h1>Cântări Creștine Videoproiector Download</h1>
    <h2>Cântările au fost preluate de pe <a href="https://cantari-crestine.com/" target="_blank">cantari-crestine.com</a></h2>
    <h2><a href="cantari-crestine-videoproiector.html">Vezi lista completă a cântărilor</a></h2>
    <div class="downloads">
        {chr(10).join(links_html)}
    </div>
</body>
</html>'''
    
    with open(os.path.join(web_dir, 'index.html'), 'w', encoding='utf-8') as f:
        f.write(html_content)

def create_zip_archive(source_dir, zip_name):
    with zipfile.ZipFile(zip_name, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, _, files in os.walk(source_dir):
            for file in files:
                if file.endswith('.pptx'):  # Only include PPTX files
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, source_dir)
                    zipf.write(file_path, arcname)

def main():
    print("Starting PowerPoint conversion process...")
    
    # Clean up existing presentations
    pptx_dir = os.path.join('data', 'pptx')
    # if os.path.exists(pptx_dir):
    #     print("Cleaning up existing presentations...")
    #     for root, dirs, files in os.walk(pptx_dir, topdown=False):
    #         for name in files:
    #             os.remove(os.path.join(root, name))
    #         for name in dirs:
    #             os.rmdir(os.path.join(root, name))
    #     os.rmdir(pptx_dir)
    #
    # # Process songs
    # print("Loading songs data...")
    songs = load_songs()
    # total_songs = len(songs)
    # print(f"Converting {total_songs} songs to PowerPoint...")
    # for i, song in enumerate(songs, 1):
    #     print(f"Converting song {i}/{total_songs}: {song['title']}")
    #     create_pptx_for_song(song)
    #
    # # Create zip archives for each category
    # print("\nCreating zip archives for each category...")
    # for category_dir in os.listdir(pptx_dir):
    #     category_path = os.path.join(pptx_dir, category_dir)
    #     if os.path.isdir(category_path):
    #         print(f"Creating zip for {category_dir}...")
    #         zip_name = os.path.join(pptx_dir, f"{category_dir}.zip")
    #         create_zip_archive(category_path, zip_name)
    #
    # # Create one zip with all songs
    # print("\nCreating zip with all songs...")
    # create_zip_archive(pptx_dir, os.path.join(pptx_dir, 'Toate cantarile crestine videoproiector.zip'))

    web_dir = os.path.join('data', 'web', 'static')
    os.mkdir(web_dir)

    # Generate index.html with download links
    print("\nGenerating index.html...")
    zip_files = [f for f in os.listdir(pptx_dir) if f.endswith('.zip')]
    generate_index_html(web_dir, zip_files)
    # Generate detailed HTML page with all songs
    print("\nGenerating detailed HTML page...")
    generate_detailed_html(web_dir, songs)
    print("Done!")

if __name__ == '__main__':
    main()
